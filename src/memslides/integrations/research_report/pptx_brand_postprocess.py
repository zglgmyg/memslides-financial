"""Standalone SJTU styling for an already-generated financial PPTX.

This module deliberately does not participate in DeckDesigner, manuscript,
HTML, memory, or export orchestration. It reads a completed PPTX and writes a
separate branded copy while preserving all existing content and geometry.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import tempfile
from collections import Counter
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


SJTU_RED = "A62038"
SJTU_LIGHT_GRAY = "BFBFBF"
SJTU_CHAMPAGNE = "E0CFBD"
SJTU_LOGO_NAME = "SJTU Official Seal"

_DEFAULT_LOGO = Path(__file__).parent / "assets" / "sjtu" / "sjtu-seal-white.png"
_LEGACY_TITLE_COLORS = {"1E3A5F", SJTU_RED}
_COLOR_MAP = {
    "1E3A5F": SJTU_RED,
    "2563EB": SJTU_RED,
    "1E40AF": SJTU_RED,
    "D97706": SJTU_RED,
    "F59E0B": SJTU_RED,
    "CBD5E1": SJTU_LIGHT_GRAY,
    "DBEAFE": SJTU_CHAMPAGNE,
    "93C5FD": SJTU_CHAMPAGNE,
}
_DARK_BACKGROUND_COLOR_MAP = {
    **_COLOR_MAP,
    "2563EB": SJTU_CHAMPAGNE,
    "1E40AF": SJTU_CHAMPAGNE,
    "D97706": SJTU_CHAMPAGNE,
    "F59E0B": SJTU_CHAMPAGNE,
    "94A3B8": SJTU_CHAMPAGNE,
    "64748B": SJTU_CHAMPAGNE,
    "334155": SJTU_LIGHT_GRAY,
}


class SJTUBrandPostprocessError(RuntimeError):
    """Raised when a safe, structure-preserving postprocess is not possible."""


@dataclass(frozen=True)
class SJTUBrandPostprocessReport:
    input_pptx: Path
    output_pptx: Path
    slide_count: int
    content_slides: tuple[int, ...]
    color_replacements: dict[str, int]
    logos_added: int
    logos_preserved: int

    def as_dict(self) -> dict[str, Any]:
        payload = asdict(self)
        payload["input_pptx"] = str(self.input_pptx)
        payload["output_pptx"] = str(self.output_pptx)
        payload["content_slides"] = list(self.content_slides)
        return payload


def _sha256_bytes(payload: bytes) -> str:
    return hashlib.sha256(payload).hexdigest()


def _rgb_hex(color: Any) -> str | None:
    try:
        if color.type == MSO_COLOR_TYPE.RGB:
            return str(color.rgb).upper()
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _replace_color(
    color: Any, counts: Counter[str], color_map: dict[str, str]
) -> None:
    source = _rgb_hex(color)
    target = color_map.get(source or "")
    if source is None or target is None or source == target:
        return
    color.rgb = RGBColor.from_string(target)
    counts[f"#{source}->#{target}"] += 1


def _replace_fill(fill: Any, counts: Counter[str], color_map: dict[str, str]) -> None:
    try:
        if fill.type == MSO_FILL_TYPE.SOLID:
            _replace_color(fill.fore_color, counts, color_map)
    except (AttributeError, TypeError, ValueError):
        pass


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _replace_text_frame(
    text_frame: Any, counts: Counter[str], color_map: dict[str, str]
) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            _replace_color(run.font.color, counts, color_map)


def _replace_shape_colors(
    shape: Any, counts: Counter[str], color_map: dict[str, str]
) -> None:
    if hasattr(shape, "fill"):
        _replace_fill(shape.fill, counts, color_map)
    try:
        _replace_fill(shape.line.fill, counts, color_map)
    except (AttributeError, TypeError, ValueError):
        pass
    if getattr(shape, "has_text_frame", False):
        _replace_text_frame(shape.text_frame, counts, color_map)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                _replace_fill(cell.fill, counts, color_map)
                _replace_text_frame(cell.text_frame, counts, color_map)


def _solid_fill_hex(shape: Any) -> str | None:
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            return _rgb_hex(shape.fill.fore_color)
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _find_title_bar(slide: Any, slide_width: int) -> Any | None:
    tolerance = Inches(0.03)
    minimum_height = Inches(0.45)
    maximum_height = Inches(0.80)
    for shape in _iter_shapes(slide.shapes):
        if (
            abs(int(shape.left)) <= tolerance
            and abs(int(shape.top)) <= tolerance
            and int(shape.width) >= int(slide_width * 0.95)
            and minimum_height <= int(shape.height) <= maximum_height
            and _solid_fill_hex(shape) in _LEGACY_TITLE_COLORS
        ):
            return shape
    return None


def _has_full_legacy_dark_background(
    slide: Any, slide_width: int, slide_height: int
) -> bool:
    for shape in _iter_shapes(slide.shapes):
        if (
            int(shape.left) <= Inches(0.03)
            and int(shape.top) <= Inches(0.03)
            and int(shape.width) >= int(slide_width * 0.95)
            and int(shape.height) >= int(slide_height * 0.95)
            and _solid_fill_hex(shape) == "1E3A5F"
        ):
            return True
    return False


def _outline_content_slides(outline_path: Path, expected_count: int) -> tuple[int, ...]:
    try:
        payload = json.loads(outline_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise SJTUBrandPostprocessError(f"Unable to read outline: {outline_path}") from exc
    slides = payload.get("slides") if isinstance(payload, dict) else None
    if not isinstance(slides, list) or len(slides) != expected_count:
        raise SJTUBrandPostprocessError(
            f"Outline/PPTX slide count mismatch: outline={len(slides or [])}, pptx={expected_count}"
        )
    return tuple(
        index
        for index, slide in enumerate(slides, start=1)
        if isinstance(slide, dict) and str(slide.get("page_role", "")) == "content"
    )


def _shape_rect(shape: Any) -> tuple[int, int, int, int]:
    return (
        int(shape.left),
        int(shape.top),
        int(shape.left + shape.width),
        int(shape.top + shape.height),
    )


def _rects_overlap(left: tuple[int, int, int, int], right: tuple[int, int, int, int]) -> bool:
    return not (
        left[2] <= right[0]
        or right[2] <= left[0]
        or left[3] <= right[1]
        or right[3] <= left[1]
    )


def _assert_logo_space(slide: Any, title_bar: Any, logo_rect: tuple[int, int, int, int]) -> None:
    for shape in _iter_shapes(slide.shapes):
        if shape.shape_id == title_bar.shape_id or shape.name == SJTU_LOGO_NAME:
            continue
        rect = _shape_rect(shape)
        is_full_slide_background = (
            rect[0] == 0
            and rect[1] == 0
            and rect[2] >= int(title_bar.width * 0.95)
            and rect[3] > int(title_bar.height * 4)
        )
        if not is_full_slide_background and _rects_overlap(rect, logo_rect):
            raise SJTUBrandPostprocessError(
                f"Logo would overlap existing shape '{shape.name}' on a content slide."
            )


def _text_snapshot(prs: Any) -> tuple[tuple[int, int, str], ...]:
    snapshot: list[tuple[int, int, str]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide.shapes):
            if shape.name == SJTU_LOGO_NAME:
                continue
            if getattr(shape, "has_text_frame", False):
                snapshot.append((slide_number, int(shape.shape_id), shape.text))
            if getattr(shape, "has_table", False):
                table_text = "\n".join(
                    cell.text for row in shape.table.rows for cell in row.cells
                )
                snapshot.append((slide_number, int(shape.shape_id), table_text))
    return tuple(snapshot)


def _geometry_snapshot(prs: Any) -> tuple[tuple[int, int, int, int, int, int], ...]:
    snapshot: list[tuple[int, int, int, int, int, int]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide.shapes):
            if shape.name == SJTU_LOGO_NAME:
                continue
            snapshot.append(
                (
                    slide_number,
                    int(shape.shape_id),
                    int(shape.left),
                    int(shape.top),
                    int(shape.width),
                    int(shape.height),
                )
            )
    return tuple(snapshot)


def _picture_snapshot(prs: Any) -> tuple[tuple[int, int, str], ...]:
    snapshot: list[tuple[int, int, str]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name != SJTU_LOGO_NAME:
                snapshot.append(
                    (slide_number, int(shape.shape_id), _sha256_bytes(shape.image.blob))
                )
    return tuple(snapshot)


def _validate_saved_deck(
    *,
    path: Path,
    expected_slide_count: int,
    expected_content_slides: tuple[int, ...],
    expected_text: tuple[tuple[int, int, str], ...],
    expected_geometry: tuple[tuple[int, int, int, int, int, int], ...],
    expected_pictures: tuple[tuple[int, int, str], ...],
) -> None:
    prs = Presentation(path)
    if len(prs.slides) != expected_slide_count:
        raise SJTUBrandPostprocessError("Postprocess changed the slide count.")
    if _text_snapshot(prs) != expected_text:
        raise SJTUBrandPostprocessError("Postprocess changed existing slide text.")
    if _geometry_snapshot(prs) != expected_geometry:
        raise SJTUBrandPostprocessError("Postprocess changed existing shape geometry.")
    if _picture_snapshot(prs) != expected_pictures:
        raise SJTUBrandPostprocessError("Postprocess changed an existing image asset.")

    for slide_number in expected_content_slides:
        slide = prs.slides[slide_number - 1]
        title_bar = _find_title_bar(slide, prs.slide_width)
        if title_bar is None or _solid_fill_hex(title_bar) != SJTU_RED:
            raise SJTUBrandPostprocessError(
                f"Slide {slide_number} does not have the required SJTU-red title bar."
            )
        logo_count = sum(
            shape.name == SJTU_LOGO_NAME for shape in _iter_shapes(slide.shapes)
        )
        if logo_count != 1:
            raise SJTUBrandPostprocessError(
                f"Slide {slide_number} has {logo_count} SJTU seals; expected exactly one."
            )


def apply_sjtu_brand_to_pptx(
    *,
    input_pptx: str | Path,
    output_pptx: str | Path,
    outline_path: str | Path | None = None,
    logo_path: str | Path | None = None,
) -> SJTUBrandPostprocessReport:
    """Write an SJTU-colored copy of a completed PPTX without changing layout."""

    source = Path(input_pptx).expanduser().resolve()
    destination = Path(output_pptx).expanduser().resolve()
    logo = Path(logo_path).expanduser().resolve() if logo_path else _DEFAULT_LOGO.resolve()
    if not source.is_file() or source.suffix.lower() != ".pptx":
        raise SJTUBrandPostprocessError(f"Input PPTX does not exist: {source}")
    if source == destination:
        raise SJTUBrandPostprocessError("Input and output PPTX must be different files.")
    if not logo.is_file():
        raise SJTUBrandPostprocessError(f"SJTU logo does not exist: {logo}")

    prs = Presentation(source)
    slide_count = len(prs.slides)
    title_bars = {
        index: title_bar
        for index, slide in enumerate(prs.slides, start=1)
        if (title_bar := _find_title_bar(slide, prs.slide_width)) is not None
    }
    detected_content = tuple(title_bars)
    if outline_path is not None:
        declared_content = _outline_content_slides(Path(outline_path).resolve(), slide_count)
        missing_bars = sorted(set(declared_content) - set(detected_content))
        if missing_bars:
            raise SJTUBrandPostprocessError(
                "Outline declares content slides without an existing title bar: "
                + ", ".join(map(str, missing_bars))
            )
        content_slides = declared_content
    else:
        content_slides = detected_content

    expected_text = _text_snapshot(prs)
    expected_geometry = _geometry_snapshot(prs)
    expected_pictures = _picture_snapshot(prs)
    replacements: Counter[str] = Counter()
    dark_background_slides = {
        index
        for index, slide in enumerate(prs.slides, start=1)
        if _has_full_legacy_dark_background(slide, prs.slide_width, prs.slide_height)
    }
    for slide_number, slide in enumerate(prs.slides, start=1):
        color_map = (
            _DARK_BACKGROUND_COLOR_MAP
            if slide_number in dark_background_slides
            else _COLOR_MAP
        )
        _replace_fill(slide.background.fill, replacements, color_map)
        for shape in _iter_shapes(slide.shapes):
            _replace_shape_colors(shape, replacements, color_map)

    logo_size = int(Inches(0.36))
    right_margin = int(Inches(0.14))
    logos_added = 0
    logos_preserved = 0
    for slide_number in content_slides:
        slide = prs.slides[slide_number - 1]
        existing = [shape for shape in slide.shapes if shape.name == SJTU_LOGO_NAME]
        if existing:
            if len(existing) != 1:
                raise SJTUBrandPostprocessError(
                    f"Slide {slide_number} already has multiple SJTU seals."
                )
            logos_preserved += 1
            continue
        title_bar = title_bars[slide_number]
        left = int(prs.slide_width) - right_margin - logo_size
        top = max(0, (int(title_bar.height) - logo_size) // 2)
        logo_rect = (left, top, left + logo_size, top + logo_size)
        _assert_logo_space(slide, title_bar, logo_rect)
        picture = slide.shapes.add_picture(
            str(logo), left, top, width=logo_size, height=logo_size
        )
        picture.name = SJTU_LOGO_NAME
        logos_added += 1

    destination.parent.mkdir(parents=True, exist_ok=True)
    file_handle, temporary_name = tempfile.mkstemp(
        suffix=".pptx", prefix=".sjtu-postprocess-", dir=destination.parent
    )
    os.close(file_handle)
    temporary = Path(temporary_name)
    try:
        prs.save(temporary)
        _validate_saved_deck(
            path=temporary,
            expected_slide_count=slide_count,
            expected_content_slides=content_slides,
            expected_text=expected_text,
            expected_geometry=expected_geometry,
            expected_pictures=expected_pictures,
        )
        temporary.replace(destination)
    finally:
        temporary.unlink(missing_ok=True)

    return SJTUBrandPostprocessReport(
        input_pptx=source,
        output_pptx=destination,
        slide_count=slide_count,
        content_slides=content_slides,
        color_replacements=dict(sorted(replacements.items())),
        logos_added=logos_added,
        logos_preserved=logos_preserved,
    )


def _parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Apply SJTU colors and a fixed official seal to a completed financial PPTX."
    )
    parser.add_argument("--input", required=True, type=Path, dest="input_pptx")
    parser.add_argument("--output", required=True, type=Path, dest="output_pptx")
    parser.add_argument(
        "--outline",
        type=Path,
        help="Optional slide_outline.json; fails if a declared content page lacks a title bar.",
    )
    parser.add_argument("--logo", type=Path, help="Optional transparent square PNG override")
    return parser


def main() -> int:
    args = _parser().parse_args()
    report = apply_sjtu_brand_to_pptx(
        input_pptx=args.input_pptx,
        output_pptx=args.output_pptx,
        outline_path=args.outline,
        logo_path=args.logo,
    )
    print(json.dumps(report.as_dict(), ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

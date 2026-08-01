from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from memslides.integrations.research_report.pptx_brand_postprocess import (
    SJTUBrandPostprocessError,
    SJTU_CHAMPAGNE,
    SJTU_LIGHT_GRAY,
    SJTU_LOGO_NAME,
    SJTU_RED,
    _solid_fill_hex,
    apply_sjtu_brand_to_pptx,
)


def _fixture_pptx(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    stripe = cover.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor.from_string("1E3A5F")

    content = prs.slides.add_slide(blank)
    title_bar = content.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.62)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor.from_string("1E3A5F")
    title = content.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(8), Inches(0.3))
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Original content title"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor.from_string("FFFFFF")
    accent = content.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6),
        Inches(1.2),
        Inches(2),
        Inches(1),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string("DBEAFE")
    accent.line.color.rgb = RGBColor.from_string("CBD5E1")
    accent.text = "Keep geometry and text"
    accent.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string("2563EB")

    ending = prs.slides.add_slide(blank)
    background = ending.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor.from_string("1E3A5F")
    ending_text = ending.shapes.add_textbox(
        Inches(3), Inches(3), Inches(7), Inches(1)
    )
    ending_text.text = "Existing final page"
    ending_text.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(
        "F59E0B"
    )

    prs.save(path)


def _shape_named(slide, name: str):
    return [shape for shape in slide.shapes if shape.name == name]


def test_postprocess_recolors_native_shapes_and_adds_only_content_logo(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "sjtu.pptx"
    _fixture_pptx(source)

    report = apply_sjtu_brand_to_pptx(input_pptx=source, output_pptx=output)
    branded = Presentation(output)

    assert report.slide_count == 3
    assert report.content_slides == (2,)
    assert report.logos_added == 1
    assert report.logos_preserved == 0
    assert _solid_fill_hex(branded.slides[0].shapes[0]) == SJTU_RED
    assert _solid_fill_hex(branded.slides[1].shapes[0]) == SJTU_RED
    assert _solid_fill_hex(branded.slides[1].shapes[2]) == SJTU_CHAMPAGNE
    assert str(branded.slides[1].shapes[2].line.color.rgb) == SJTU_LIGHT_GRAY
    assert (
        str(branded.slides[1].shapes[2].text_frame.paragraphs[0].runs[0].font.color.rgb)
        == SJTU_RED
    )
    assert not _shape_named(branded.slides[0], SJTU_LOGO_NAME)
    assert len(_shape_named(branded.slides[1], SJTU_LOGO_NAME)) == 1
    assert not _shape_named(branded.slides[2], SJTU_LOGO_NAME)
    assert branded.slides[1].shapes[1].text == "Original content title"
    assert branded.slides[1].shapes[2].text == "Keep geometry and text"
    assert branded.slides[2].shapes[1].text == "Existing final page"
    assert (
        str(branded.slides[2].shapes[1].text_frame.paragraphs[0].runs[0].font.color.rgb)
        == SJTU_CHAMPAGNE
    )


def test_postprocess_is_idempotent(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    first = tmp_path / "first.pptx"
    second = tmp_path / "second.pptx"
    _fixture_pptx(source)

    apply_sjtu_brand_to_pptx(input_pptx=source, output_pptx=first)
    report = apply_sjtu_brand_to_pptx(input_pptx=first, output_pptx=second)

    assert report.logos_added == 0
    assert report.logos_preserved == 1
    assert len(_shape_named(Presentation(second).slides[1], SJTU_LOGO_NAME)) == 1


def test_outline_mismatch_fails_instead_of_styling_wrong_page(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "sjtu.pptx"
    outline = tmp_path / "slide_outline.json"
    _fixture_pptx(source)
    outline.write_text(
        json.dumps(
            {
                "slides": [
                    {"page_role": "title"},
                    {"page_role": "content"},
                    {"page_role": "content"},
                ]
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(SJTUBrandPostprocessError, match="without an existing title bar: 3"):
        apply_sjtu_brand_to_pptx(
            input_pptx=source, output_pptx=output, outline_path=outline
        )

    assert not output.exists()


def test_postprocess_refuses_to_overwrite_source(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    _fixture_pptx(source)

    with pytest.raises(SJTUBrandPostprocessError, match="must be different"):
        apply_sjtu_brand_to_pptx(input_pptx=source, output_pptx=source)
